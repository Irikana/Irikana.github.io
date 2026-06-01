import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'pylib'))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree
from playwright.sync_api import sync_playwright

HTML_PATH = os.path.join(os.path.dirname(__file__), 'symmetry-operations-presentation.html')
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), 'symmetry-operations-presentation.pptx')
SCREENSHOT_DIR = os.path.join(os.path.dirname(__file__), '_screenshots')

SLIDE_W = 13.333
SLIDE_H = 7.5

os.makedirs(SCREENSHOT_DIR, exist_ok=True)

def add_slide_transition(slide, spd='med'):
    existing = slide._element.find(qn('p:transition'))
    if existing is not None:
        slide._element.remove(existing)
    transition = etree.SubElement(slide._element, qn('p:transition'))
    transition.set('spd', spd)
    transition.set('advClick', '1')
    fade = etree.SubElement(transition, qn('p:fade'))
    fade.set('thruBlk', '0')

SCALE_CSS = '''
var style = document.createElement('style');
style.textContent = `
  .slide { padding: 28px 50px !important; }
  .st { font-size: 52px !important; }
  .ssub { font-size: 22px !important; }
  .sc h2 { font-size: 32px !important; margin-bottom: 18px !important; }
  .sc p { font-size: 21px !important; line-height: 1.8 !important; }
  .cbox { font-size: 21px !important; padding: 18px 28px !important; }
  .fbox { font-size: 20px !important; padding: 16px 24px !important; }
  .cbox2 h4 { font-size: 17px !important; }
  .cbox2 p { font-size: 17px !important; line-height: 1.7 !important; }
  .ch { font-size: 20px !important; padding: 12px !important; }
  .sflow .stx { font-size: 19px !important; }
  .sflow .sn { width: 30px !important; height: 30px !important; font-size: 14px !important; }
  .flow .fb { font-size: 19px !important; padding: 9px 18px !important; }
  .flow .ar { font-size: 26px !important; }
  .choices { max-width: 640px !important; gap: 14px !important; }
  .dcon svg { transform: scale(1.25) !important; transform-origin: center !important; }
  .en-dim { font-size: 17px !important; }
`;
document.head.appendChild(style);
'''

def capture_slides():
    shots = []

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={'width': 1920, 'height': 1080},
            device_scale_factor=2
        )
        page.goto('file:///' + HTML_PATH.replace('\\', '/'))
        page.wait_for_timeout(800)

        page.evaluate('''() => {
            var nav = document.querySelector('.nav');
            if (nav) nav.style.display = 'none';
        }''')

        page.evaluate(SCALE_CSS)
        page.wait_for_timeout(300)

        def max_step(idx):
            return page.evaluate('''(idx) => {
                var m = 0;
                var el = document.querySelectorAll('.slide')[idx];
                if (!el) return 0;
                el.querySelectorAll('[data-step]').forEach(function(e) {
                    var s = parseInt(e.getAttribute('data-step'));
                    if (s > m) m = s;
                });
                return m;
            }''', idx)

        def go_to_state(idx, step_val):
            page.evaluate('''([idx, stepVal]) => {
                var slides = document.querySelectorAll('.slide');
                slides.forEach(function(s, i) {
                    if (i === idx) s.classList.add('active');
                    else s.classList.remove('active');
                });
                var el = slides[idx];
                if (!el) return;
                el.querySelectorAll('[data-step]').forEach(function(e) {
                    var ns = parseInt(e.getAttribute('data-step'));
                    if (ns <= stepVal) e.classList.add('vis');
                    else e.classList.remove('vis');
                });
            }''', [idx, step_val])

        def trigger_special_logic(idx, step_val):
            page.evaluate('''([idx, stepVal]) => {
                if (idx === 2) {
                    var c = document.getElementById('chC');
                    if (c) { if (stepVal >= 2) c.classList.add('ok'); else c.classList.remove('ok'); }
                }
                if (idx === 3) {
                    var f = document.getElementById('chF');
                    if (f) { if (stepVal >= 2) f.classList.add('ok'); else f.classList.remove('ok'); }
                }
                if (idx === 4) {
                    var app = document.getElementById('apparatus');
                    var gso = document.getElementById('greenSolidOrig');
                    var gbo = document.getElementById('greenBallOrig');
                    var gss = document.getElementById('greenSolidShift');
                    var gbs = document.getElementById('greenBallShift');
                    var gd = document.getElementById('greenDashed');
                    var bd = document.getElementById('blueDashed');
                    if (app) {
                        if (stepVal >= 1 && stepVal <= 3) app.style.transform = 'rotate(90deg)';
                        else app.style.transform = 'rotate(0deg)';
                    }
                    var showMirror = (stepVal >= 2 && stepVal <= 3);
                    if (gso) gso.style.opacity = showMirror ? '0' : '1';
                    if (gbo) gbo.style.opacity = showMirror ? '0' : '1';
                    if (gss) gss.style.opacity = showMirror ? '1' : '0';
                    if (gbs) gbs.style.opacity = showMirror ? '1' : '0';
                    if (gd) gd.style.opacity = showMirror ? '1' : '0';
                    if (bd) bd.style.opacity = showMirror ? '1' : '0';
                }
            }''', [idx, step_val])

        total_slides = page.evaluate('document.querySelectorAll(".slide").length')
        print(f'Found {total_slides} slides in HTML')

        for i in range(total_slides):
            ms = max_step(i)
            print(f'  Slide {i}: max_step={ms}')
            for s in range(ms + 1):
                go_to_state(i, s)
                trigger_special_logic(i, s)
                page.wait_for_timeout(700)
                fname = f'slide_{i:02d}_step_{s:02d}.png'
                fpath = os.path.join(SCREENSHOT_DIR, fname)
                page.screenshot(path=fpath, full_page=False)
                shots.append((fpath, i, s, ms))
                print(f'    Captured step {s}/{ms}')

        browser.close()

    return shots

def build_pptx(shots):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    blank_layout = prs.slide_layouts[6]

    for idx, (img_path, slide_idx, step_val, max_steps) in enumerate(shots):
        slide = prs.slides.add_slide(blank_layout)

        pic = slide.shapes.add_picture(
            img_path,
            Inches(0), Inches(0),
            Inches(SLIDE_W), Inches(SLIDE_H)
        )

        add_slide_transition(slide, spd='med')

    return prs

if __name__ == '__main__':
    print('Capturing slides from HTML...')
    shots = capture_slides()
    print(f'Captured {len(shots)} screenshots total')

    print('Building PPTX...')
    prs = build_pptx(shots)
    prs.save(OUTPUT_PATH)
    print(f'Saved to {OUTPUT_PATH}')
    print(f'Total slides in PPTX: {len(prs.slides)}')
